"""
Document parser — extracts text from PDF, DOCX, XLSX, and other formats.
Used by file_read (auto-detect) and the attach system.
"""

import json
from pathlib import Path

# Supported extensions and their parsers
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".doc", ".pptx"}


def can_parse(path: str) -> bool:
    return Path(path).suffix.lower() in SUPPORTED_EXTENSIONS


def parse_document(path: str, max_chars: int = 50000) -> dict:
    """Extract text from a document file.

    Returns {"content": str, "pages": int, "format": str}
    or {"error": str} on failure.
    """
    p = Path(path)
    ext = p.suffix.lower()

    if not p.exists():
        return {"error": f"File not found: {path}"}

    try:
        if ext == ".pdf":
            return _parse_pdf(p, max_chars)
        elif ext == ".docx":
            return _parse_docx(p, max_chars)
        elif ext in (".xlsx", ".xls"):
            return _parse_xlsx(p, max_chars)
        elif ext == ".pptx":
            return _parse_pptx(p, max_chars)
        elif ext == ".doc":
            return {"error": "Legacy .doc format not supported. Convert to .docx first."}
        else:
            return {"error": f"Unsupported format: {ext}"}
    except Exception as e:
        return {"error": f"Failed to parse {p.name}: {type(e).__name__}: {e}"}


def _parse_pdf(path: Path, max_chars: int) -> dict:
    import fitz  # PyMuPDF
    doc = fitz.open(str(path))
    pages = []
    total = 0
    for i, page in enumerate(doc):
        text = page.get_text()
        if total + len(text) > max_chars:
            text = text[:max_chars - total]
            pages.append(f"--- Page {i + 1} ---\n{text}\n...(truncated)")
            total = max_chars
            break
        pages.append(f"--- Page {i + 1} ---\n{text}")
        total += len(text)
    doc.close()
    return {
        "content": "\n".join(pages),
        "pages": len(doc) if hasattr(doc, '__len__') else len(pages),
        "format": "pdf",
    }


def _parse_docx(path: Path, max_chars: int) -> dict:
    from docx import Document
    doc = Document(str(path))
    parts = []
    total = 0
    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            continue
        if total + len(text) > max_chars:
            parts.append(text[:max_chars - total] + "\n...(truncated)")
            break
        parts.append(text)
        total += len(text)
    return {
        "content": "\n".join(parts),
        "pages": 1,  # docx doesn't have a clean page count
        "format": "docx",
    }


def _parse_xlsx(path: Path, max_chars: int) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    total = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells)
            if total + len(line) > max_chars:
                parts.append("...(truncated)")
                total = max_chars
                break
            parts.append(line)
            total += len(line)
        if total >= max_chars:
            break
    wb.close()
    return {
        "content": "\n".join(parts),
        "pages": len(wb.sheetnames),
        "format": "xlsx",
    }


def _parse_pptx(path: Path, max_chars: int) -> dict:
    # Basic extraction without python-pptx dependency
    # Falls back to zip-based extraction of slide text
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        total = 0
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texts.append(para.text)
            text = "\n".join(texts)
            if total + len(text) > max_chars:
                parts.append(f"--- Slide {i + 1} ---\n{text[:max_chars - total]}\n...(truncated)")
                break
            parts.append(f"--- Slide {i + 1} ---\n{text}")
            total += len(text)
        return {"content": "\n".join(parts), "pages": len(prs.slides), "format": "pptx"}
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}
